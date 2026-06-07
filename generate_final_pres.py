"""
HealthHubQA FYP — 5-slide presentation (CLEAN version)
Title on own page. No GPT-4 (invalid). Less wordy. No GenAI mention.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colours
DARK_BG      = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE  = RGBColor(0x00, 0x9F, 0xD4)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_RED   = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_ORANGE= RGBColor(0xF3, 0x9C, 0x12)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY     = RGBColor(0x88, 0x88, 0x99)
DARK_TEXT     = RGBColor(0x2C, 0x3E, 0x50)
TABLE_HEADER = RGBColor(0x00, 0x7B, 0xA7)
TABLE_ALT    = RGBColor(0xF0, 0xF4, 0xF8)
TABLE_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def text(slide, l, t, w, h, txt, sz=18, bold=False, col=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt; p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = col; p.font.name = 'Calibri'; p.alignment = align
    return tf

def bullets(slide, l, t, w, h, items, sz=15, col=LIGHT_GRAY):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        r = p.add_run(); r.text = f"•  {item}"
        r.font.size = Pt(sz); r.font.color.rgb = col; r.font.name = 'Calibri'

def box(slide, l, t, w, h, fill, txt, sz=14, col=WHITE, bold=False):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = col
    p.font.bold = bold; p.font.name = 'Calibri'

def line(slide, l, t, w):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(0.04))
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT_BLUE; s.line.fill.background()

def arrow(slide, l, t, w=0.4):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(l), Inches(t), Inches(w), Inches(0.3))
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT_BLUE; s.line.fill.background()

def table(slide, l, t, w, h, data, cw=None):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w), Inches(h))
    tbl = ts.table
    if cw:
        for i, v in enumerate(cw): tbl.columns[i].width = Inches(v)
    for ri, rd in enumerate(data):
        for ci, ct in enumerate(rd):
            c = tbl.cell(ri, ci); c.text = str(ct)
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(13); p.font.name = 'Calibri'; p.alignment = PP_ALIGN.CENTER
                p.font.bold = ri == 0; p.font.color.rgb = WHITE if ri == 0 else DARK_TEXT
            c.fill.solid()
            c.fill.fore_color.rgb = TABLE_HEADER if ri == 0 else (TABLE_ALT if ri % 2 == 0 else TABLE_WHITE)
    return tbl

def notes(slide, txt):
    slide.notes_slide.notes_text_frame.text = txt


# ════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE ONLY — clean, professional
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)

text(s, 1.5, 1.5, 10.3, 1.2,
    "Investigating and Mitigating Hallucinations in\nLarge Language Models for Healthcare QA",
    sz=34, bold=True, col=WHITE, align=PP_ALIGN.CENTER)

line(s, 4.5, 2.85, 4.3)

text(s, 1.5, 3.1, 10.3, 0.8, "HealthHubQA",
    sz=48, bold=True, col=ACCENT_BLUE, align=PP_ALIGN.CENTER)

text(s, 1.5, 4.3, 10.3, 1.5,
    "Ayomide Ogun-Ajala  |  21309903\n"
    "BSc Computer Science & Software Engineering\n\n"
    "Supervisor: Dr. Kolawole Adebayo\n"
    "Maynooth University  |  2025–2026",
    sz=20, col=LIGHT_GRAY, align=PP_ALIGN.CENTER)

notes(s,
    "Say: 'My project investigates how to detect and reduce hallucinations "
    "in AI systems used for healthcare question answering. "
    "I built a system called HealthHubQA.'")


# ════════════════════════════════════════════════════════════════
# SLIDE 2: PROBLEM + RAG + PIPELINE  (no waffle)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)

text(s, 0.6, 0.3, 10, 0.5, "Problem, Approach & Pipeline", sz=26, bold=True)
line(s, 0.6, 0.85, 3)

# Problem — compact
text(s, 0.6, 1.1, 5.5, 0.4, "The Problem", sz=20, bold=True, col=ACCENT_BLUE)
bullets(s, 0.6, 1.5, 5.5, 1.2, [
    "LLMs produce confident but factually wrong answers (hallucinations)",
    "In healthcare, wrong answers can cause patient harm",
    "Can we ground answers in real medical evidence to reduce this?",
], sz=14)

# RAG — two compact boxes
text(s, 6.8, 1.1, 5.5, 0.4, "Solution: RAG (Retrieval-Augmented Generation)", sz=17, bold=True, col=ACCENT_BLUE)

box(s, 6.8, 1.55, 5.5, 0.55, RGBColor(0x5C, 0x1A, 0x1A),
    "Without RAG = Closed-book exam\nAI answers from memory — may fabricate facts", sz=12, col=WHITE)
box(s, 6.8, 2.2, 5.5, 0.65, RGBColor(0x1A, 0x3C, 0x1A),
    "With RAG = Open-book exam\nAI retrieves real medical papers first, reads them, THEN answers", sz=12, col=WHITE)

# Pipeline
text(s, 0.6, 3.2, 10, 0.4, "HealthHubQA Pipeline", sz=20, bold=True, col=ACCENT_BLUE)
pipes = [
    (0.6,  "Medical\nQuestion",       RGBColor(0x34, 0x49, 0x5E)),
    (2.85, "Retriever\n(MiniLM-L6-v2)", ACCENT_BLUE),
    (5.1,  "Top-3\nContexts",         RGBColor(0x27, 0xAE, 0x60)),
    (7.35, "Generator\n(LLM)",        RGBColor(0x8E, 0x44, 0xAD)),
    (9.6,  "Yes / No /\nMaybe",       RGBColor(0xE6, 0x7E, 0x22)),
    (11.4, "Evaluate",                RGBColor(0xC0, 0x39, 0x2B)),
]
for x, lbl, c in pipes:
    box(s, x, 3.7, 1.9, 0.75, c, lbl, sz=12, col=WHITE, bold=True)
for i in range(5):
    arrow(s, pipes[i][0] + 1.95, 3.92, w=0.35)

# Models table — compact
text(s, 0.6, 4.8, 10, 0.4, "4 Models Tested (× 2 conditions = 8 experiments)", sz=17, bold=True, col=ACCENT_BLUE)
md = [
    ["Model", "Type", "Why Chosen", "Params"],
    ["Flan-T5", "Instruction-tuned (Google)", "Free baseline, follows instructions", "250M"],
    ["BioBART", "Biomedical pre-trained", "Tests domain-specific knowledge", "140M"],
    ["DistilBART", "Distilled summariser", "Tests weaker architecture", "306M"],
    ["Phi-3 Mini", "Small LLM (Microsoft)", "Tests compact local model", "3.8B"],
]
table(s, 0.6, 5.2, 10, 2.0, md, cw=[1.8, 2.2, 3.5, 1.0])

# Dataset info
box(s, 11.0, 5.2, 1.7, 0.9, RGBColor(0x22, 0x22, 0x3A),
    "PubMedQA\n500 docs\n100 test Qs", sz=12, col=LIGHT_GRAY)
box(s, 11.0, 6.2, 1.7, 0.6, RGBColor(0x22, 0x22, 0x3A),
    "Phi-3: 10 Qs\n(CPU limit)", sz=11, col=MED_GRAY)

notes(s,
    "Say: 'LLMs hallucinate — they make things up. In healthcare that is dangerous.\n\n"
    "My solution is RAG. Think of it as giving a student an open-book exam instead of closed-book. "
    "The AI retrieves the top 3 most relevant medical papers, reads them, then answers.\n\n"
    "Here is the pipeline: question comes in, retriever finds relevant papers using sentence embeddings, "
    "passes them to the AI, which outputs yes/no/maybe, and we evaluate against ground truth.\n\n"
    "I tested 4 models: Flan-T5 as a baseline, BioBART for medical domain knowledge, "
    "DistilBART to test a weaker architecture, and Phi-3 Mini as a compact local option.\n\n"
    "All tested on 100 PubMedQA medical questions (10 for Phi-3 due to CPU constraints).'")


# ════════════════════════════════════════════════════════════════
# SLIDE 3: RESULTS TABLE  (the star slide)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)

text(s, 0.6, 0.3, 10, 0.5, "Results: RAG Impact on Accuracy", sz=26, bold=True)
line(s, 0.6, 0.85, 3)

# Results table — 4 valid models only
rd = [
    ["Model", "No-RAG", "RAG", "Δ Change", "F1 (RAG)", "Precision", "Recall", "N"],
    ["Flan-T5",    "40%", "55%", "+15%  ↑", "0.237", "0.189", "0.316", "100"],
    ["BioBART",    "58%", "58%", "  0%  —", "0.245", "0.193", "0.333", "100"],
    ["DistilBART", "39%", "32%", " -7%  ↓", "0.271", "0.285", "0.272", "100"],
    ["Phi-3 Mini", "30%", "40%", "+10%  ↑", "0.277", "0.310", "0.317", "10"],
]
tbl = table(s, 0.6, 1.1, 12.1, 2.2, rd, cw=[1.7, 1.3, 1.3, 1.5, 1.3, 1.3, 1.3, 0.8])

# Colour delta cells
colours = [ACCENT_GREEN, ACCENT_ORANGE, ACCENT_RED, ACCENT_GREEN]
for ri, c in enumerate(colours, start=1):
    cell = tbl.cell(ri, 3)
    for p in cell.text_frame.paragraphs:
        p.font.color.rgb = c; p.font.bold = True

# 3 interpretation boxes
box(s, 0.6, 3.6, 3.8, 1.5, RGBColor(0x1A, 0x3C, 0x1A),
    "Flan-T5: +15% with RAG\n\nInstruction-tuned model\nknows how to USE the\nretrieved context",
    sz=13, col=ACCENT_GREEN)

box(s, 4.7, 3.6, 3.8, 1.5, RGBColor(0x2C, 0x2C, 0x15),
    "BioBART: 58% both ways\n\nMedical pre-training already\nbakes in domain knowledge —\nRAG adds nothing new",
    sz=13, col=ACCENT_ORANGE)

box(s, 8.8, 3.6, 3.9, 1.5, RGBColor(0x3C, 0x1A, 0x1A),
    "DistilBART: -7% with RAG\n\nSummarisation model gets\nCONFUSED by extra context —\nRAG is NOT universal",
    sz=13, col=ACCENT_RED)

# Bottom clarification
text(s, 0.6, 5.4, 12, 0.8,
    "Accuracy = correct yes/no/maybe predictions  |  Macro F1 = balanced across all 3 classes  |  "
    "Phi-3 ran on CPU (10 samples only — promising signal, needs more data to confirm)  |  "
    "BioBART: best no-RAG model  |  Flan-T5: biggest reliable RAG improvement",
    sz=12, col=MED_GRAY)

notes(s,
    "Say: 'This is my key slide — the results.\n\n"
    "Flan-T5 went from 40% to 55% with RAG — a 15-point improvement. "
    "It is instruction-tuned, meaning it was trained to follow instructions. "
    "So when you give it medical papers and say answer based on these, it knows what to do.\n\n"
    "BioBART stayed at 58% both with and without RAG. "
    "This model was already pre-trained on biomedical literature — "
    "it already has medical knowledge baked in. RAG does not add new information.\n\n"
    "DistilBART actually DROPPED from 39% to 32%. "
    "This is my most interesting finding. DistilBART is a summarisation model — "
    "when you give it extra medical context, it tries to summarise instead of answering. "
    "The extra text confused it. This shows RAG is NOT a one-size-fits-all solution.\n\n"
    "Phi-3 improved from 30% to 40% but only on 10 samples due to CPU limits. "
    "Promising but needs more testing.\n\n"
    "The take-away: RAG helps models trained to follow instructions, "
    "but not all architectures can handle extra context.'")


# ════════════════════════════════════════════════════════════════
# SLIDE 4: KEY FINDINGS + CHALLENGES
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)

text(s, 0.6, 0.3, 10, 0.5, "Key Findings & Engineering Challenges", sz=26, bold=True)
line(s, 0.6, 0.85, 3)

# 3 findings
box(s, 0.6, 1.1, 6.0, 0.9, RGBColor(0x1A, 0x3C, 0x1A),
    "1. RAG Improves Instruction-Tuned Models\nFlan-T5: 40% → 55%  (+15 points)",
    sz=14, col=ACCENT_GREEN, bold=True)

box(s, 6.9, 1.1, 5.8, 0.9, RGBColor(0x3C, 0x1A, 0x1A),
    "2. RAG Can Hurt Wrong Architectures\nDistilBART: 39% → 32%  (-7 points)",
    sz=14, col=ACCENT_RED, bold=True)

box(s, 0.6, 2.2, 6.0, 0.9, RGBColor(0x2C, 0x2C, 0x15),
    "3. Domain Pre-Training = Strong Alternative\nBioBART: 58% without any retrieval",
    sz=14, col=ACCENT_ORANGE, bold=True)

box(s, 6.9, 2.2, 5.8, 0.9, RGBColor(0x1B, 0x2A, 0x3C),
    "4. Practical Constraints Matter\nRAM: Phi-3 (2.2GB) vs Llama (4.6GB — too large for laptop)",
    sz=14, col=ACCENT_BLUE, bold=True)

# Challenges
text(s, 0.6, 3.4, 10, 0.4, "Engineering Challenges Solved", sz=18, bold=True)
ch = [
    ["Challenge", "Cause", "Solution", "Result"],
    ["FAISS crash", "NumPy 2.x broke API", "sklearn cosine similarity", "Same results, stable"],
    ["Llama 3 too large", "4.6GB RAM needed", "Switched to Phi-3 Mini (2.2GB)", "Runs on laptop"],
    ["T5 tokenizer broke", "Library version conflict", "Replaced with BioBART", "Better: biomedical model"],
    ["API quota exhausted", "No billing credits", "Error handling + logging", "No data lost"],
]
table(s, 0.6, 3.85, 12.1, 2.0, ch, cw=[2.2, 2.5, 3.5, 2.8])

# Literature strip at bottom
text(s, 0.6, 6.2, 12, 0.5,
    "Literature: Lewis et al. (2020) RAG  |  Zuo & Jiang (2025) MedHallBench  |  "
    "Ji et al. (2023) Self-Reflection  |  Asgari et al. (2025) Clinical Safety  |  "
    "Sok et al. (2025) MetaRAG",
    sz=11, col=MED_GRAY)

notes(s,
    "Say: 'My key findings:\n\n"
    "First — RAG helps instruction-tuned models. Flan-T5 gained 15 points. "
    "These models are trained to follow instructions, so they use the retrieved context properly.\n\n"
    "Second — RAG can actually hurt the wrong architecture. DistilBART dropped 7 points. "
    "It is a summarisation model that gets confused by extra context.\n\n"
    "Third — domain pre-training is a strong alternative. BioBART reached 58% without any retrieval "
    "because medical knowledge is already in its weights.\n\n"
    "Fourth — practical constraints matter. Llama 3 needed 4.6GB RAM which my laptop could not handle, "
    "so I switched to Phi-3 Mini at 2.2GB.\n\n"
    "I also solved real engineering problems — FAISS crashed, a tokenizer broke, API credits ran out. "
    "Each challenge led to a better solution.\n\n"
    "My work builds on Lewis et al. who introduced RAG, and Zuo & Jiang for medical hallucination benchmarking.'")


# ════════════════════════════════════════════════════════════════
# SLIDE 5: CONCLUSION + FUTURE WORK + THANK YOU
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)

text(s, 0.6, 0.3, 10, 0.5, "Conclusion & Future Work", sz=26, bold=True)
line(s, 0.6, 0.85, 3)

# 3 takeaways
box(s, 0.6, 1.1, 12.1, 0.6, RGBColor(0x1A, 0x3C, 0x1A),
    "1.  RAG reduces hallucinations for instruction-tuned models  (Flan-T5: 40% → 55%)",
    sz=16, col=ACCENT_GREEN, bold=True)
box(s, 0.6, 1.85, 12.1, 0.6, RGBColor(0x3C, 0x1A, 0x1A),
    "2.  Not all models benefit — architecture matters  (DistilBART: 39% → 32%  ↓)",
    sz=16, col=ACCENT_RED, bold=True)
box(s, 0.6, 2.6, 12.1, 0.6, RGBColor(0x2C, 0x2C, 0x15),
    "3.  Domain pre-training (BioBART 58%) is a viable alternative to retrieval augmentation",
    sz=16, col=ACCENT_ORANGE, bold=True)

# Contribution
text(s, 0.8, 3.4, 11.5, 0.7,
    "Contribution: A comparative multi-model study proving RAG effectiveness\n"
    "varies by architecture — it is not a universal solution for healthcare QA.",
    sz=16, bold=True, col=WHITE, align=PP_ALIGN.CENTER)

# Future work
text(s, 0.6, 4.2, 10, 0.3, "Future Work", sz=16, bold=True, col=ACCENT_BLUE)
fw = [
    (0.6,  "Self-Reflection\nAI checks own answer\nbefore output", ACCENT_BLUE),
    (2.8,  "More Datasets\nMedQA, MedMCQA\ncross-domain", RGBColor(0x27, 0xAE, 0x60)),
    (5.0,  "FACTSCORE\nDirect hallucination\nmeasurement", RGBColor(0x8E, 0x44, 0xAD)),
    (7.2,  "Web Demo\nStreamlit/Flask\nlive QA interface", RGBColor(0xE6, 0x7E, 0x22)),
    (9.4,  "GPU Testing\nLarger models at\nfull scale", RGBColor(0xC0, 0x39, 0x2B)),
    (11.6, "Clinical Safety\nRisk-aware\nevaluation", RGBColor(0x16, 0xA0, 0x85)),
]
for x, txt, c in fw:
    box(s, x, 4.55, 2.0, 0.9, c, txt, sz=10, col=WHITE)

# Thank you
text(s, 0.6, 5.7, 12, 0.5, "Thank You — Questions?",
    sz=28, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
text(s, 0.6, 6.3, 12, 0.4,
    "4 models  •  8 experiments  •  500 documents  •  17 scripts  •  2,500+ lines of code",
    sz=13, col=MED_GRAY, align=PP_ALIGN.CENTER)
text(s, 0.6, 6.7, 12, 0.4,
    "Ayomide Ogun-Ajala  |  21309903  |  Dr. Kolawole Adebayo  |  Maynooth University",
    sz=12, col=MED_GRAY, align=PP_ALIGN.CENTER)

notes(s,
    "Say: 'Three conclusions:\n\n"
    "First — RAG works for instruction-tuned models. Flan-T5 gained 15 points.\n"
    "Second — RAG is not universal. DistilBART got worse.\n"
    "Third — domain pre-training like BioBART is a strong alternative.\n\n"
    "My contribution: proving that RAG effectiveness depends on model architecture.\n\n"
    "Future extensions include self-reflection loops, more datasets, "
    "direct hallucination metrics, a web demo, GPU-scale testing, and clinical safety scoring.\n\n"
    "Thank you. Happy to take questions.'")


# Save
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'HealthHubQA_FYP_FINAL.pptx')
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
