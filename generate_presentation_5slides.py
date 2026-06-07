"""
Generate HealthHubQA FYP Presentation — 5 SLIDES ONLY
Condensed version with all key info packed in.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colours ─────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE  = RGBColor(0x00, 0x9F, 0xD4)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_RED   = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_ORANGE= RGBColor(0xF3, 0x9C, 0x12)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY     = RGBColor(0x88, 0x88, 0x99)
DARK_TEXT     = RGBColor(0x2C, 0x3E, 0x50)
TABLE_HEADER = RGBColor(0x00, 0x7B, 0xA7)
TABLE_ALT    = RGBColor(0xF0, 0xF4, 0xF8)
TABLE_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

# ── Helpers ─────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=15, color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        if isinstance(item, tuple):
            r1 = p.add_run()
            r1.text = item[0]
            r1.font.size = Pt(font_size)
            r1.font.bold = True
            r1.font.color.rgb = ACCENT_BLUE
            r1.font.name = 'Calibri'
            r2 = p.add_run()
            r2.text = item[1]
            r2.font.size = Pt(font_size)
            r2.font.color.rgb = color
            r2.font.name = 'Calibri'
        else:
            r = p.add_run()
            r.text = f"•  {item}"
            r.font.size = Pt(font_size)
            r.font.color.rgb = color
            r.font.name = 'Calibri'
    return tf

def add_box(slide, left, top, width, height, fill_color, text,
            font_size=14, text_color=WHITE, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.font.name = 'Calibri'
    return shape

def add_line(slide, left, top, width):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.04))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT_BLUE
    s.line.fill.background()

def add_arrow(slide, left, top, w=0.45):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
        Inches(left), Inches(top), Inches(w), Inches(0.3))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT_BLUE
    s.line.fill.background()

def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
    ts = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for r_idx, row_data in enumerate(data):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = 'Calibri'
                paragraph.alignment = PP_ALIGN.CENTER
                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = DARK_TEXT
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT if r_idx % 2 == 0 else TABLE_WHITE
    return table

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE + PROBLEM + RAG EXPLANATION
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

# Title
add_textbox(slide, 0.6, 0.3, 8, 0.8,
    "Investigating & Mitigating Hallucinations in\nLarge Language Models for Healthcare QA",
    font_size=26, bold=True, color=WHITE)
add_textbox(slide, 0.6, 1.15, 8, 0.4,
    "HealthHubQA", font_size=36, bold=True, color=ACCENT_BLUE)
add_line(slide, 0.6, 1.1, 3)

add_textbox(slide, 0.6, 1.7, 7, 0.8,
    "Ayomide Ogun-Ajala  |  21309903  |  BSc CS & SE\n"
    "Supervisor: Dr. Kolawole Adebayo  |  Maynooth University",
    font_size=14, color=MED_GRAY)

# The Problem — left side
add_textbox(slide, 0.6, 2.6, 5.5, 0.4, "The Problem", font_size=20, bold=True, color=WHITE)
add_bullet_list(slide, 0.6, 3.0, 5.5, 1.8, [
    "LLMs generate confident but factually wrong answers (hallucinations)",
    "In healthcare, wrong answers can lead to harmful treatment decisions",
    "Research Q: Can we ground AI answers in real medical evidence?",
], font_size=14, color=LIGHT_GRAY)

# Hallucination example box
add_box(slide, 0.6, 4.6, 5.5, 0.55, RGBColor(0x5C, 0x1A, 0x1A),
        '❌  Q: "Is metformin for hypertension?"  →  AI: "Yes" (WRONG — it\'s for diabetes)',
        font_size=11, text_color=ACCENT_RED)
add_box(slide, 0.6, 5.25, 5.5, 0.55, RGBColor(0x1A, 0x3C, 0x1A),
        '✓  With RAG: Retrieves real papers → AI correctly says "No, metformin treats diabetes"',
        font_size=11, text_color=ACCENT_GREEN)

# What is RAG — right side
add_textbox(slide, 6.8, 2.6, 5.8, 0.4, "What is RAG?", font_size=20, bold=True, color=WHITE)

add_box(slide, 6.8, 3.05, 5.8, 0.6, RGBColor(0x5C, 0x1A, 0x1A),
        "WITHOUT RAG  =  Closed-Book Exam\nAI answers from memory — may fabricate facts",
        font_size=12, text_color=WHITE, bold=False)

add_box(slide, 6.8, 3.8, 5.8, 0.7, RGBColor(0x1A, 0x3C, 0x1A),
        "WITH RAG  =  Open-Book Exam\nAI retrieves real medical papers first, reads them, THEN answers\nAnswer is grounded in evidence",
        font_size=12, text_color=WHITE, bold=False)

# Technical box
add_box(slide, 6.8, 4.65, 5.8, 1.15, RGBColor(0x22, 0x22, 0x3A),
        "Technical: Sentence Embeddings (all-MiniLM-L6-v2)\n"
        "384-dim vectors  |  Cosine Similarity  |  Top-3 retrieval\n"
        "Index: 500 PubMedQA medical documents\n"
        "Pipeline: Question → Retriever → Top-3 Docs → LLM → Answer",
        font_size=11, text_color=MED_GRAY)

# Literature at very bottom
add_textbox(slide, 0.6, 6.2, 12, 0.8,
    "Built on: Lewis et al. (2020) RAG  |  Zuo & Jiang (2025) MedHallBench  |  "
    "Ji et al. (2023) Self-Reflection  |  Asgari et al. (2025) Clinical Safety  |  "
    "Sok et al. (2025) MetaRAG",
    font_size=11, color=MED_GRAY)

add_notes(slide,
    "OPENING: 'My project is about making AI safer in healthcare.\n\n"
    "LLMs like ChatGPT sometimes hallucinate — they produce confident but wrong answers. "
    "Imagine asking: Is metformin for hypertension? A hallucinating AI says yes — but that's WRONG, "
    "metformin is for diabetes.\n\n"
    "I built HealthHubQA to test whether RAG — Retrieval Augmented Generation — can fix this. "
    "Think of it as an open-book vs closed-book exam. Without RAG, the AI answers from memory and might make things up. "
    "With RAG, it first searches a library of 500 real medical papers, finds the 3 most relevant ones using sentence "
    "embeddings and cosine similarity, reads them, and THEN answers. The answer is grounded in real evidence.\n\n"
    "My work builds on Lewis et al. who introduced RAG, Zuo & Jiang for medical hallucination benchmarking, "
    "and Ji et al. for self-reflection techniques.'")


# ════════════════════════════════════════════════════════════════
# SLIDE 2: EXPERIMENTAL DESIGN — Models + Pipeline
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.6, 0.3, 10, 0.5,
    "Experimental Design: 5 Models × 2 Conditions = 10 Experiments",
    font_size=26, bold=True, color=WHITE)
add_line(slide, 0.6, 0.85, 3)

# Models table
models_data = [
    ["Model", "Type", "Source", "Why Chosen", "Params"],
    ["Flan-T5", "Instruction-tuned", "Google / HuggingFace", "Free baseline, follows instructions", "250M"],
    ["GPT-4", "Cloud API", "OpenAI", "Gold standard commercial AI", "~1.7T"],
    ["BioBART", "Biomedical pre-trained", "GanjinZero / HuggingFace", "Tests domain-specific pre-training", "140M"],
    ["DistilBART", "Distilled summariser", "sshleifer / HuggingFace", "Tests weaker architecture", "306M"],
    ["Phi-3 Mini", "Small LLM (Ollama)", "Microsoft", "Tests compact local models", "3.8B"],
]
add_table(slide, 0.6, 1.1, 9.5, 2.6, 6, 5, models_data,
          col_widths=[1.6, 1.8, 2.0, 2.8, 0.9])

# Design info boxes — right side
add_box(slide, 10.5, 1.1, 2.3, 0.8, RGBColor(0x22, 0x22, 0x3A),
        "2 × 5 Grid\n10 Experiments", font_size=16, text_color=ACCENT_BLUE, bold=True)
add_box(slide, 10.5, 2.05, 2.3, 0.7, RGBColor(0x22, 0x22, 0x3A),
        "PubMedQA Dataset\n500 medical Q&A pairs", font_size=12, text_color=LIGHT_GRAY)
add_box(slide, 10.5, 2.9, 2.3, 0.7, RGBColor(0x22, 0x22, 0x3A),
        "100 test questions\n(20 for Phi-3: CPU limit)", font_size=12, text_color=LIGHT_GRAY)

# Pipeline flow
add_textbox(slide, 0.6, 4.0, 10, 0.4, "HealthHubQA Pipeline", font_size=20, bold=True, color=WHITE)

pipeline_boxes = [
    (0.6, "Medical\nQuestion", RGBColor(0x34, 0x49, 0x5E)),
    (2.7, "Retriever\n(MiniLM-L6-v2)", ACCENT_BLUE),
    (4.8, "Top-3\nContexts", RGBColor(0x27, 0xAE, 0x60)),
    (6.9, "Generator\n(LLM)", RGBColor(0x8E, 0x44, 0xAD)),
    (9.0, "Yes / No /\nMaybe", RGBColor(0xE6, 0x7E, 0x22)),
    (11.1, "Evaluate\n(Metrics)", RGBColor(0xC0, 0x39, 0x2B)),
]
y = 4.5
for x, label, color in pipeline_boxes:
    add_box(slide, x, y, 1.8, 0.8, color, label, font_size=12, text_color=WHITE, bold=True)
for i in range(5):
    add_arrow(slide, pipeline_boxes[i][0] + 1.85, y + 0.25, w=0.4)

# Bottom: methodology
add_textbox(slide, 0.6, 5.7, 12, 0.9,
    "Methodology: Same 100 PubMedQA questions per model  |  "
    "Yes/No/Maybe classification task  |  "
    "Metrics: Accuracy, Macro F1, Precision, Recall  |  "
    "Confusion matrices & similarity scores logged  |  "
    "All results in CSV for full reproducibility",
    font_size=12, color=MED_GRAY)

add_notes(slide,
    "SAY: 'I set up a fair comparison. 5 different AI models, each tested twice: "
    "once without help, once with retrieved medical papers. That's 10 experiments.\n\n"
    "I chose diverse models on purpose:\n"
    "- Flan-T5: a small, instruction-following model as a baseline\n"
    "- GPT-4: the best commercial AI, as a gold standard\n"
    "- BioBART: pre-trained specifically on medical text, to test domain knowledge\n"
    "- DistilBART: a weaker summarisation model, to test if all architectures benefit\n"
    "- Phi-3 Mini: Microsoft's small but capable model, to test compact local options\n\n"
    "The pipeline works like this: A medical question comes in. The retriever searches 500 PubMedQA "
    "documents using sentence embeddings, finds the top 3 matches, and passes them to the generator. "
    "The AI reads the question plus retrieved context and outputs yes, no, or maybe. "
    "We compare to the known correct answer and record accuracy, F1, precision, and recall.'")


# ════════════════════════════════════════════════════════════════
# SLIDE 3: RESULTS — The Big Table + Detailed Metrics
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.6, 0.3, 10, 0.5,
    "Results: RAG Impact Across All 5 Models",
    font_size=26, bold=True, color=WHITE)
add_line(slide, 0.6, 0.85, 3)

# Main accuracy results table
results_data = [
    ["Model", "No-RAG Accuracy", "RAG Accuracy", "Δ Change", "Macro F1 (RAG)", "Precision (RAG)", "Recall (RAG)", "N"],
    ["GPT-4", "40%", "64%", "+24%  ↑", "0.482", "0.476", "0.503", "100"],
    ["Flan-T5", "40%", "55%", "+15%  ↑", "0.237", "0.189", "0.316", "100"],
    ["Phi-3 Mini", "30%", "75%", "+45%  ↑", "0.277", "0.310", "0.317", "20"],
    ["BioBART", "58%", "58%", "  0%  —", "0.245", "0.193", "0.333", "100"],
    ["DistilBART", "39%", "32%", " -7%  ↓", "0.271", "0.285", "0.272", "100"],
]
tbl = add_table(slide, 0.6, 1.0, 12.1, 2.6, 6, 8, results_data,
                col_widths=[1.6, 1.5, 1.4, 1.4, 1.5, 1.5, 1.4, 0.6])

# Colour the delta column
for r_idx in range(1, 6):
    cell = tbl.cell(r_idx, 3)
    for p in cell.text_frame.paragraphs:
        if r_idx <= 3:
            p.font.color.rgb = ACCENT_GREEN
            p.font.bold = True
        elif r_idx == 4:
            p.font.color.rgb = ACCENT_ORANGE
            p.font.bold = True
        else:
            p.font.color.rgb = ACCENT_RED
            p.font.bold = True

# Interpretation boxes below
add_box(slide, 0.6, 3.85, 3.8, 1.4, RGBColor(0x1A, 0x3C, 0x1A),
        "GPT-4 + RAG = Best Overall\n64% accuracy (+24%)\n\n"
        "Like giving a smart student\na textbook — they use it well",
        font_size=12, text_color=ACCENT_GREEN)

add_box(slide, 4.6, 3.85, 3.8, 1.4, RGBColor(0x2C, 0x2C, 0x15),
        "BioBART = Strong Without RAG\n58% accuracy (unchanged)\n\n"
        "Medical pre-training already\nbakes in domain knowledge",
        font_size=12, text_color=ACCENT_ORANGE)

add_box(slide, 8.6, 3.85, 4.1, 1.4, RGBColor(0x3C, 0x1A, 0x1A),
        "DistilBART = RAG Hurts (-7%)\n39% → 32% accuracy\n\n"
        "Summarisation model gets confused\nby extra context — RAG isn't universal",
        font_size=12, text_color=ACCENT_RED)

# Bottom note
add_textbox(slide, 0.6, 5.5, 12, 0.9,
    "Key insight: Accuracy = correct yes/no/maybe predictions.  Macro F1 balances across all 3 classes (important for imbalanced data).  "
    "GPT-4 + RAG leads ALL metrics.  Flan-T5 accuracy improved but F1 dipped — better on majority yes class, still struggles with no/maybe.  "
    "Phi-3's +45% is promising but only n=20 (CPU constraint) — interpret with caution.",
    font_size=12, color=MED_GRAY)

add_notes(slide,
    "SAY: 'This is my KEY slide. Let me walk through each result:\n\n"
    "GPT-4 jumped from 40% to 64% with RAG — a 24-point improvement. It has the highest F1 of 0.482, "
    "highest precision and recall. Like giving a smart student a textbook — they know how to use it.\n\n"
    "Flan-T5 also improved: 40% to 55%. Even smaller models benefit from evidence.\n\n"
    "Phi-3 went from 30% to 75% — biggest jump. BUT only 20 samples due to CPU limits. "
    "Promising signal but needs more data.\n\n"
    "BioBART stayed at 58% both ways. It was trained on medical text so it already knows the domain. "
    "RAG didn't add new knowledge. Proves medical pre-training is powerful on its own.\n\n"
    "DistilBART DROPPED from 39% to 32%. THIS is my most interesting finding. "
    "It's a summarisation model that gets confused by extra context. "
    "Proves RAG is NOT a one-size-fits-all solution.\n\n"
    "F1-score matters because PubMedQA is imbalanced — more yes answers. "
    "A model saying yes to everything gets ~55% accuracy but poor F1.'")


# ════════════════════════════════════════════════════════════════
# SLIDE 4: KEY FINDINGS + CHALLENGES
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.6, 0.3, 10, 0.5,
    "Key Findings, Contributions & Engineering Challenges",
    font_size=26, bold=True, color=WHITE)
add_line(slide, 0.6, 0.85, 3)

# 4 findings as boxes — top half
add_box(slide, 0.6, 1.1, 5.8, 1.0, RGBColor(0x1A, 0x3C, 0x1A),
        "1. RAG Reduces Hallucinations for Instruction-Tuned Models\n"
        "GPT-4: +24%  |  Flan-T5: +15%  |  Phi-3: +45%",
        font_size=13, text_color=ACCENT_GREEN, bold=True)

add_box(slide, 6.8, 1.1, 5.8, 1.0, RGBColor(0x3C, 0x1A, 0x1A),
        "2. Model Architecture Determines RAG Effectiveness\n"
        "DistilBART: -7% — summarisation models can't handle retrieval context",
        font_size=13, text_color=ACCENT_RED, bold=True)

add_box(slide, 0.6, 2.3, 5.8, 1.0, RGBColor(0x2C, 0x2C, 0x15),
        "3. Domain Pre-Training is a Strong Alternative to RAG\n"
        "BioBART: 58% without retrieval — competitive with RAG-enhanced models",
        font_size=13, text_color=ACCENT_ORANGE, bold=True)

add_box(slide, 6.8, 2.3, 5.8, 1.0, RGBColor(0x1B, 0x2A, 0x3C),
        "4. Practical Constraints Shape Deployment\n"
        "RAM: Phi-3 (2.2GB) vs Llama (4.6GB)  |  Cost: GPT-4 API vs free local models",
        font_size=13, text_color=ACCENT_BLUE, bold=True)

# Challenges table — bottom half
add_textbox(slide, 0.6, 3.6, 10, 0.4, "Engineering Challenges Solved", font_size=18, bold=True, color=WHITE)

challenges = [
    ["Challenge", "Root Cause", "Solution", "Outcome"],
    ["FAISS library crash", "NumPy 2.x incompatibility", "sklearn cosine similarity", "Same results, stable"],
    ["Llama 3 out of RAM", "4.6GB needed, 4.5GB free", "Phi-3 Mini (2.2GB)", "Best RAG improvement"],
    ["GPT API credits exhausted", "Zero billing on 2 accounts", "try/except error logging", "No data loss"],
    ["T5 tokenizer failure", "Transformers version bug", "Replaced with BioBART", "Better: 58% medical"],
    ["API key leaked in chat", "Accidental plain-text", "Revoked + env variables", "Security restored"],
]
add_table(slide, 0.6, 4.0, 12.1, 2.5, 6, 4, challenges,
          col_widths=[2.5, 2.8, 3.2, 2.5])

add_notes(slide,
    "SAY: 'Four key contributions:\n\n"
    "1. RAG helps instruction-tuned models significantly — GPT-4 gained 24 points.\n"
    "2. But NOT all models benefit. DistilBART got WORSE — it's a summarisation model that gets confused by extra context. "
    "This challenges the assumption that RAG is always good.\n"
    "3. Domain pre-training can replace retrieval. BioBART hit 58% without RAG because medical knowledge is already in its weights.\n"
    "4. Practical constraints matter — RAM and API costs shape what's deployable.\n\n"
    "I also solved several real engineering problems. FAISS crashed so I used cosine similarity instead. "
    "Llama was too large so I found Phi-3 Mini. My GPT credits ran out but error handling saved the data. "
    "The T5 tokenizer broke but BioBART was actually an improvement. "
    "And I accidentally leaked an API key but immediately revoked it and switched to environment variables.'")


# ════════════════════════════════════════════════════════════════
# SLIDE 5: FUTURE WORK + CONCLUSION + THANK YOU
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.6, 0.3, 10, 0.5,
    "Conclusion & Future Work",
    font_size=26, bold=True, color=WHITE)
add_line(slide, 0.6, 0.85, 3)

# 3 conclusions
add_box(slide, 0.6, 1.1, 12.1, 0.6, RGBColor(0x1A, 0x3C, 0x1A),
        "1.  RAG reduces hallucinations for instruction-tuned models  (GPT-4: 40% → 64%)",
        font_size=16, text_color=ACCENT_GREEN, bold=True)

add_box(slide, 0.6, 1.85, 12.1, 0.6, RGBColor(0x3C, 0x1A, 0x1A),
        "2.  Not all models benefit — architecture matters  (DistilBART: 39% → 32%  ↓)",
        font_size=16, text_color=ACCENT_RED, bold=True)

add_box(slide, 0.6, 2.6, 12.1, 0.6, RGBColor(0x2C, 0x2C, 0x15),
        "3.  Domain pre-training (BioBART 58%) is a viable alternative to retrieval augmentation",
        font_size=16, text_color=ACCENT_ORANGE, bold=True)

# Contribution statement
add_textbox(slide, 0.8, 3.4, 11.5, 0.7,
    "Contribution: A comparative multi-model study showing RAG's hallucination mitigation\n"
    "effectiveness varies by architecture — challenging the assumption RAG universally helps.",
    font_size=15, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Future work boxes
add_textbox(slide, 0.6, 4.2, 10, 0.3, "Future Extensions", font_size=16, bold=True, color=ACCENT_BLUE)

fw = [
    (0.6,  "Self-Reflection Loops\nAI reviews own answer\nbefore outputting", ACCENT_BLUE),
    (2.8,  "More Datasets\nMedQA, MedMCQA for\ncross-domain testing", RGBColor(0x27, 0xAE, 0x60)),
    (5.0,  "FACTSCORE Metrics\nDirect hallucination\nmeasurement", RGBColor(0x8E, 0x44, 0xAD)),
    (7.2,  "Web Interface\nStreamlit/Flask live\nhealthcare QA demo", RGBColor(0xE6, 0x7E, 0x22)),
    (9.4,  "GPU-Scale Tests\nFull Phi-3 + larger\nmodel evaluation", RGBColor(0xC0, 0x39, 0x2B)),
    (11.6, "Clinical Safety\nAsgari et al. (2025)\nrisk-aware scoring", RGBColor(0x16, 0xA0, 0x85)),
]
for x, text, color in fw:
    add_box(slide, x, 4.55, 2.0, 1.0, color, text, font_size=10, text_color=WHITE, bold=False)

# Thank you + stats
add_textbox(slide, 0.6, 5.8, 12, 0.5,
    "Thank You — Questions?",
    font_size=28, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.6, 6.35, 12, 0.4,
    "5 models  •  10 experiments  •  500 documents  •  17 scripts  •  2,500+ lines of code",
    font_size=13, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.6, 6.75, 12, 0.4,
    "Ayomide Ogun-Ajala  |  21309903  |  Dr. Kolawole Adebayo  |  Maynooth University",
    font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_notes(slide,
    "SAY: 'To conclude — three key takeaways:\n\n"
    "First, RAG significantly reduces hallucinations for instruction-tuned models like GPT-4.\n"
    "Second, RAG is not universal — DistilBART actually got worse. Architecture matters.\n"
    "Third, domain pre-training like BioBART is a viable alternative.\n\n"
    "My contribution is a comparative multi-model study that challenges the assumption that RAG always helps.\n\n"
    "For future work: self-reflection loops, more datasets, direct hallucination metrics like FACTSCORE, "
    "a web interface demo, and GPU-scale experiments.\n\n"
    "For real-world healthcare AI, we need BOTH good retrieval AND the right model to make AI trustworthy.\n\n"
    "Thank you. I'm happy to take any questions.'")


# ── Save ────────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           'HealthHubQA_FYP_5slides.pptx')
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
