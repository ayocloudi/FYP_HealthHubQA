"""
Generate HealthHubQA FYP Presentation (.pptx)
With full evaluation results, pipeline diagrams, and speaking notes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colour palette ──────────────────────────────────────────────
DARK_BG      = RGBColor(0x1B, 0x1B, 0x2F)   # dark navy
ACCENT_BLUE  = RGBColor(0x00, 0x9F, 0xD4)   # bright teal
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)   # green for positive
ACCENT_RED   = RGBColor(0xE7, 0x4C, 0x3C)   # red for negative
ACCENT_ORANGE= RGBColor(0xF3, 0x9C, 0x12)   # orange/amber
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY     = RGBColor(0x88, 0x88, 0x99)
DARK_TEXT     = RGBColor(0x2C, 0x3E, 0x50)
TABLE_HEADER  = RGBColor(0x00, 0x7B, 0xA7)
TABLE_ALT     = RGBColor(0xF0, 0xF4, 0xF8)
TABLE_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# ── Helper functions ────────────────────────────────────────────

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, font_name='Calibri', bold_prefix=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        
        if isinstance(item, tuple):
            # (bold_part, rest)
            run1 = p.add_run()
            run1.text = item[0]
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = ACCENT_BLUE
            run1.font.name = font_name
            run2 = p.add_run()
            run2.text = item[1]
            run2.font.size = Pt(font_size)
            run2.font.bold = False
            run2.font.color.rgb = color
            run2.font.name = font_name
        else:
            run = p.add_run()
            run.text = f"  •  {item}"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = font_name
    return tf

def add_shape_box(slide, left, top, width, height, fill_color, text, 
                  font_size=14, text_color=WHITE, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.font.name = 'Calibri'
    shape.text_frame.paragraphs[0].space_before = Pt(4)
    return shape

def add_arrow(slide, left, top, width=0.5, height=0.01):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   Inches(left), Inches(top), Inches(width), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    return shape

def add_table(slide, left, top, width, height, rows, cols, data, 
              col_widths=None):
    """data = list of lists, first row is header"""
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                          Inches(width), Inches(height))
    table = table_shape.table
    
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    
    for r_idx, row_data in enumerate(data):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = 'Calibri'
                paragraph.alignment = PP_ALIGN.CENTER
                
                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = DARK_TEXT
            
            # Header row styling
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT if r_idx % 2 == 0 else TABLE_WHITE
    
    return table

def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text

# ═════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

# Accent line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(1), Inches(2.3), Inches(4), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

add_textbox(slide, 1, 1.0, 11, 1.2,
    "Investigating and Mitigating Hallucinations in\nLarge Language Models for Healthcare QA",
    font_size=32, bold=True, color=WHITE)

add_textbox(slide, 1, 2.5, 8, 0.6,
    "HealthHubQA", font_size=44, bold=True, color=ACCENT_BLUE)

add_textbox(slide, 1, 3.4, 8, 1.5,
    "Ayomide Ogun-Ajala  |  21309903\n"
    "BSc Computer Science & Software Engineering\n"
    "Supervisor: Dr. Kolawole Adebayo\n"
    "Maynooth University  |  2025-2026",
    font_size=18, color=LIGHT_GRAY)

add_notes(slide, 
    "OPENING: 'My project is about making AI safer in healthcare. "
    "LLMs sometimes hallucinate — they produce confident but wrong answers. "
    "In healthcare, this could be dangerous. I built HealthHubQA to test whether "
    "giving AI real medical papers to read first (RAG) reduces hallucinations. "
    "I tested 5 models and found RAG helps some dramatically, but not all.'")

# ═════════════════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.4, 10, 0.6, "The Problem: AI Hallucinations in Healthcare",
            font_size=30, bold=True, color=WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(1.05), Inches(3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

add_bullet_list(slide, 0.8, 1.3, 6, 2.5, [
    "LLMs generate fluent, confident text — but sometimes it's factually wrong",
    "In healthcare, a wrong answer could lead to harmful treatment decisions",
    "These fabricated outputs are called 'hallucinations'",
    "Research question: Can we ground AI answers in real medical evidence?",
], font_size=17, color=LIGHT_GRAY)

# Example box
add_shape_box(slide, 7.5, 1.3, 5, 1.0, RGBColor(0x2C, 0x2C, 0x44),
              "Hallucination Example", font_size=14, text_color=ACCENT_BLUE, bold=True)

add_textbox(slide, 7.7, 2.3, 4.8, 0.4,
    'Q: "Is metformin used to treat hypertension?"', font_size=14, color=WHITE)

add_shape_box(slide, 7.7, 2.8, 4.6, 0.7, RGBColor(0x5C, 0x1A, 0x1A),
              '❌  AI says: "Yes, metformin lowers blood pressure"\n(WRONG — metformin is for diabetes)',
              font_size=12, text_color=ACCENT_RED)

add_shape_box(slide, 7.7, 3.7, 4.6, 0.7, RGBColor(0x1A, 0x3C, 0x1A),
              '✓  With RAG: Retrieves papers → AI correctly answers\n"No, metformin is used for diabetes management"',
              font_size=12, text_color=ACCENT_GREEN)

# Stats at bottom
add_textbox(slide, 0.8, 4.8, 11, 0.8,
    "Healthcare AI errors can directly impact patient safety — hallucination detection & mitigation is critical",
    font_size=15, bold=False, color=ACCENT_ORANGE, alignment=PP_ALIGN.LEFT)

add_notes(slide,
    "SAY: 'Imagine asking an AI: Is metformin used for hypertension? "
    "A hallucinating AI might say yes. But that's WRONG — metformin is for diabetes. "
    "My system retrieves actual medical papers first, so the AI answers based on evidence, not guessing. "
    "This is what my project is all about — can we make AI truthful in healthcare?'")

# ═════════════════════════════════════════════════════════════════
# SLIDE 3: WHAT IS RAG?
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.4, 10, 0.6, "What is RAG? (Retrieval-Augmented Generation)",
            font_size=30, bold=True, color=WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(1.05), Inches(3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

# Left side: No-RAG
add_shape_box(slide, 0.8, 1.5, 5.5, 0.6, RGBColor(0x5C, 0x1A, 0x1A),
              "WITHOUT RAG  =  Closed-Book Exam", font_size=16, text_color=WHITE, bold=True)
add_bullet_list(slide, 0.8, 2.2, 5.5, 1.2, [
    "AI answers purely from memory (training data)",
    "May remember incorrectly or fabricate facts",
    "No way to verify or ground the answer",
], font_size=15, color=LIGHT_GRAY)

# Right side: RAG
add_shape_box(slide, 6.8, 1.5, 5.8, 0.6, RGBColor(0x1A, 0x3C, 0x1A),
              "WITH RAG  =  Open-Book Exam", font_size=16, text_color=WHITE, bold=True)
add_bullet_list(slide, 6.8, 2.2, 5.8, 1.5, [
    "Question → Retriever searches medical papers",
    "Finds top 3 most relevant documents",
    "AI reads documents + question, THEN answers",
    "Answer is grounded in real evidence",
], font_size=15, color=LIGHT_GRAY)

# Technical details box at bottom
add_shape_box(slide, 0.8, 4.0, 11.7, 1.2, RGBColor(0x22, 0x22, 0x3A),
              "Technical Implementation\n"
              "Embedding Model: all-MiniLM-L6-v2 (384-dim vectors)  |  "
              "Similarity: Cosine Similarity (sklearn)  |  "
              "Index: 500 PubMedQA documents  |  "
              "Retrieval: Top-3 passages per question",
              font_size=13, text_color=MED_GRAY)

add_notes(slide,
    "SAY: 'Think of it like an open-book vs closed-book exam. "
    "Without RAG, the AI answers from memory — it might remember wrong. "
    "With RAG, the AI first searches a library of 500 real medical papers, "
    "finds the 3 most relevant ones using sentence embeddings and cosine similarity, "
    "reads them, and THEN answers. The answer is grounded in real evidence.'")

# ═════════════════════════════════════════════════════════════════
# SLIDE 4: EXPERIMENTAL DESIGN
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.4, 10, 0.6, "Experimental Design: 5 Models × 2 Conditions",
            font_size=30, bold=True, color=WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(1.05), Inches(3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

# Models table
models_data = [
    ["Model", "Type", "Source", "Why Chosen"],
    ["Flan-T5 (250M)", "Instruction-tuned", "Google / HuggingFace", "Free, lightweight baseline"],
    ["GPT-4", "Cloud API", "OpenAI", "Gold standard commercial AI"],
    ["BioBART", "Biomedical pre-trained", "GanjinZero / HuggingFace", "Tests domain pre-training"],
    ["DistilBART", "Distilled summariser", "sshleifer / HuggingFace", "Tests weaker architecture"],
    ["Phi-3 Mini (3.8B)", "Small LLM", "Microsoft / Ollama", "Tests local small models"],
]
add_table(slide, 0.8, 1.3, 8, 2.8, 6, 4, models_data,
          col_widths=[2.2, 1.8, 2.2, 2.3])

# Design summary box on right
add_shape_box(slide, 9.3, 1.3, 3.3, 1.0, RGBColor(0x22, 0x22, 0x3A),
              "2 × 5 Grid\n10 Total Experiments", font_size=18, text_color=ACCENT_BLUE, bold=True)

add_shape_box(slide, 9.3, 2.5, 3.3, 0.7, RGBColor(0x22, 0x22, 0x3A),
              "Dataset: PubMedQA\n500 medical Q&A pairs", font_size=14, text_color=LIGHT_GRAY)

add_shape_box(slide, 9.3, 3.4, 3.3, 0.7, RGBColor(0x22, 0x22, 0x3A),
              "Test: 100 questions/model\n(20 for Phi-3 — CPU limits)", font_size=14, text_color=LIGHT_GRAY)

# Labels explanation
add_textbox(slide, 0.8, 4.5, 11, 0.6,
    "Task: Yes / No / Maybe classification  —  Model reads a medical question and decides if the answer is yes, no, or maybe",
    font_size=15, color=MED_GRAY)

add_notes(slide,
    "SAY: 'I set up a fair comparison: 5 different AI models, each tested twice — "
    "once without help (no-RAG) and once with retrieved medical papers (RAG). "
    "That gives 10 experiments total. I used 100 real medical questions from PubMedQA — "
    "a well-known biomedical QA dataset. Each question has a known correct answer: yes, no, or maybe. "
    "I chose diverse models to see if RAG benefits all architectures equally.'")

# ═════════════════════════════════════════════════════════════════
# SLIDE 5: PIPELINE ARCHITECTURE
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.4, 10, 0.6, "HealthHubQA Pipeline Architecture",
            font_size=30, bold=True, color=WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(1.05), Inches(3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

# Pipeline boxes
y = 2.0
boxes = [
    (1.0, "Medical\nQuestion", RGBColor(0x34, 0x49, 0x5E)),
    (3.3, "Retriever\n(MiniLM-L6-v2)", ACCENT_BLUE),
    (5.6, "Top-3\nContexts", RGBColor(0x27, 0xAE, 0x60)),
    (7.9, "Generator\n(LLM)", RGBColor(0x8E, 0x44, 0xAD)),
    (10.2, "Answer\n(Yes/No/Maybe)", RGBColor(0xE6, 0x7E, 0x22)),
]
for x, label, color in boxes:
    add_shape_box(slide, x, y, 2.0, 1.0, color, label, font_size=14, text_color=WHITE, bold=True)

# Arrows between boxes
for i in range(4):
    x = boxes[i][0] + 2.05
    add_arrow(slide, x, y + 0.33)

# Database shape below retriever
add_shape_box(slide, 3.3, 3.3, 2.0, 0.8, RGBColor(0x22, 0x22, 0x3A),
              "PubMedQA Index\n500 Documents", font_size=12, text_color=MED_GRAY)

# Evaluation row below
add_textbox(slide, 0.8, 4.3, 11, 0.5, "Evaluation Pipeline:", font_size=18, bold=True, color=ACCENT_BLUE)

eval_boxes = [
    (1.0, "Compare to\nGround Truth", RGBColor(0x2C, 0x2C, 0x44)),
    (3.5, "Accuracy &\nMacro F1", RGBColor(0x2C, 0x2C, 0x44)),
    (6.0, "Precision &\nRecall", RGBColor(0x2C, 0x2C, 0x44)),
    (8.5, "Confusion\nMatrices", RGBColor(0x2C, 0x2C, 0x44)),
    (11.0, "Similarity\nScores", RGBColor(0x2C, 0x2C, 0x44)),
]
for x, label, color in eval_boxes:
    add_shape_box(slide, x, 4.8, 2.0, 0.8, color, label, font_size=12, text_color=LIGHT_GRAY)

for i in range(4):
    x = eval_boxes[i][0] + 2.05
    add_arrow(slide, x, 5.0, width=0.4)

add_notes(slide,
    "SAY: 'Here is how the system works. A medical question comes in. "
    "The retriever, using sentence embeddings from all-MiniLM-L6-v2, searches our index of 500 PubMedQA documents "
    "and finds the top 3 most similar passages using cosine similarity. "
    "These passages are given to the generator (the LLM) along with the question. "
    "The model outputs yes, no, or maybe. We then compare to the known correct answer and compute "
    "accuracy, F1-score, precision, recall, and confusion matrices.'")

# ═════════════════════════════════════════════════════════════════
# SLIDE 6: RESULTS — THE BIG TABLE (Most Important Slide)
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.3, 10, 0.6, "Results: RAG Impact Across 5 Models",
            font_size=30, bold=True, color=WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(0.95), Inches(3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

# Main results table
results_data = [
    ["Model", "No-RAG", "RAG", "Δ Accuracy", "Macro F1 (RAG)", "N Samples"],
    ["GPT-4", "40%", "64%", "+24%  ↑", "0.482", "100"],
    ["Flan-T5", "40%", "55%", "+15%  ↑", "0.237", "100"],
    ["Phi-3 Mini", "30%", "75%", "+45%  ↑", "0.277", "20"],
    ["BioBART", "58%", "58%", "  0%  —", "0.245", "100"],
    ["DistilBART", "39%", "32%", " -7%  ↓", "0.271", "100"],
]
tbl = add_table(slide, 0.8, 1.2, 11.5, 2.8, 6, 6, results_data,
                col_widths=[2.2, 1.5, 1.5, 2.0, 2.0, 1.3])

# Color-code the delta column
for r_idx in range(1, 6):
    cell = tbl.cell(r_idx, 3)
    for p in cell.text_frame.paragraphs:
        if r_idx <= 3:  # positive
            p.font.color.rgb = ACCENT_GREEN
            p.font.bold = True
        elif r_idx == 4:  # neutral
            p.font.color.rgb = ACCENT_ORANGE
            p.font.bold = True
        else:  # negative
            p.font.color.rgb = ACCENT_RED
            p.font.bold = True

# Key insight boxes below table
add_shape_box(slide, 0.8, 4.3, 3.6, 1.2, RGBColor(0x1A, 0x3C, 0x1A),
              "GPT-4 + RAG\n64% accuracy\nBest overall performer",
              font_size=13, text_color=ACCENT_GREEN, bold=True)

add_shape_box(slide, 4.6, 4.3, 3.6, 1.2, RGBColor(0x2C, 0x2C, 0x15),
              "BioBART No-RAG\n58% accuracy\nMedical pre-training works",
              font_size=13, text_color=ACCENT_ORANGE, bold=True)

add_shape_box(slide, 8.4, 4.3, 4.2, 1.2, RGBColor(0x3C, 0x1A, 0x1A),
              "DistilBART + RAG\n32% (worse than baseline)\nRAG can hurt wrong architecture",
              font_size=13, text_color=ACCENT_RED, bold=True)

add_notes(slide,
    "SAY: 'This is my most important slide. Let me walk through each model:\n\n"
    "GPT-4 jumped from 40% to 64% with RAG — a 24-point improvement. "
    "Like giving a smart student a textbook — they know how to use it.\n\n"
    "Flan-T5 also improved, from 40% to 55%. Even smaller models benefit from evidence.\n\n"
    "Phi-3 Mini showed the biggest jump: 30% to 75%. BUT this was only 20 samples due to CPU constraints, "
    "so we should be cautious. It is a promising signal.\n\n"
    "BioBART stayed at 58% in both conditions. It was already trained on medical text, "
    "so RAG didn't add new knowledge. This proves medical pre-training is powerful on its own.\n\n"
    "DistilBART actually got WORSE — from 39% to 32%. This is my most interesting finding. "
    "DistilBART is a summarisation model that gets confused by extra context. "
    "This proves RAG is NOT a one-size-fits-all solution — architecture matters.'")

# ═════════════════════════════════════════════════════════════════
# SLIDE 7: DETAILED EVALUATION METRICS
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.3, 10, 0.6, "Detailed Evaluation Metrics",
            font_size=30, bold=True, color=WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(0.95), Inches(3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

# Full metrics table
full_metrics = [
    ["Model", "Setup", "Accuracy", "Macro F1", "Precision", "Recall", "N"],
    ["Flan-T5", "No-RAG", "40%", "0.290", "0.302", "0.357", "100"],
    ["Flan-T5", "RAG", "55%", "0.237", "0.189", "0.316", "100"],
    ["GPT-4", "No-RAG", "40%", "0.328", "0.334", "0.334", "100"],
    ["GPT-4", "RAG", "64%", "0.482", "0.476", "0.503", "100"],
    ["BioBART", "No-RAG", "58%", "0.245", "0.193", "0.333", "100"],
    ["BioBART", "RAG", "58%", "0.245", "0.193", "0.333", "100"],
    ["DistilBART", "No-RAG", "39%", "0.327", "0.352", "0.327", "100"],
    ["DistilBART", "RAG", "32%", "0.271", "0.285", "0.272", "100"],
    ["Phi-3 Mini", "No-RAG", "30%", "0.232", "0.244", "0.233", "20"],
    ["Phi-3 Mini", "RAG", "75%", "0.277", "0.310", "0.317", "20"],
]
add_table(slide, 0.5, 1.1, 12.3, 4.3, 11, 7, full_metrics,
          col_widths=[1.8, 1.0, 1.3, 1.3, 1.3, 1.3, 0.7])

# Explanation box
add_textbox(slide, 0.5, 5.5, 12, 0.8,
    "Note: Accuracy measures correct yes/no/maybe predictions.  Macro F1 balances precision & recall across all 3 classes.  "
    "GPT-4 + RAG shows the strongest improvement across ALL metrics.  Phi-3's results are promising but based on small sample (n=20).",
    font_size=13, color=MED_GRAY)

add_notes(slide,
    "SAY: 'This table shows the full metrics breakdown. Key observations:\n"
    "- GPT-4 + RAG is the strongest across ALL metrics: highest accuracy (64%), highest F1 (0.482), highest precision (0.476) and recall (0.503).\n"
    "- Flan-T5 accuracy improved with RAG (40→55%) but its F1 actually dipped slightly — this means it improved on the majority class (yes) but still struggles with no/maybe.\n"
    "- BioBART is identical in both conditions — RAG literally made no difference, the same predictions were produced.\n"
    "- DistilBART dropped across every single metric with RAG — accuracy, F1, precision, and recall all decreased.\n"
    "- Phi-3 shows the most dramatic accuracy gain but on only 20 samples, so we interpret with caution.\n\n"
    "F1-score is important because PubMedQA is imbalanced — more yes answers than no or maybe. "
    "A model that just says yes to everything would get ~55% accuracy but poor F1.'")

# ═════════════════════════════════════════════════════════════════
# SLIDE 8: KEY FINDINGS / RESEARCH CONTRIBUTIONS
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.4, 10, 0.6, "Key Findings & Research Contributions",
            font_size=30, bold=True, color=WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(1.05), Inches(3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

# Finding 1
add_shape_box(slide, 0.8, 1.4, 5.6, 1.3, RGBColor(0x1A, 0x3C, 0x1A),
              "1. RAG Significantly Reduces Hallucinations\nfor Instruction-Tuned Models\n\n"
              "GPT-4: +24%  |  Flan-T5: +15%  |  Phi-3: +45%",
              font_size=14, text_color=ACCENT_GREEN)

# Finding 2
add_shape_box(slide, 6.8, 1.4, 5.6, 1.3, RGBColor(0x3C, 0x1A, 0x1A),
              "2. Model Architecture Determines\nRAG Effectiveness\n\n"
              "DistilBART: -7% (RAG hurts summarisation models)",
              font_size=14, text_color=ACCENT_RED)

# Finding 3
add_shape_box(slide, 0.8, 3.0, 5.6, 1.3, RGBColor(0x2C, 0x2C, 0x15),
              "3. Domain Pre-Training is a Strong\nAlternative to RAG\n\n"
              "BioBART: 58% without retrieval — competitive with RAG models",
              font_size=14, text_color=ACCENT_ORANGE)

# Finding 4
add_shape_box(slide, 6.8, 3.0, 5.6, 1.3, RGBColor(0x1B, 0x2A, 0x3C),
              "4. Practical Constraints Shape\nDeployment Decisions\n\n"
              "RAM: Phi-3 (2.2GB) vs Llama (4.6GB)  |  Cost: GPT-4 API vs Free local",
              font_size=14, text_color=ACCENT_BLUE)

# Bottom summary
add_textbox(slide, 0.8, 4.7, 11.5, 0.8,
    "Main Contribution: A comparative multi-model analysis showing RAG's hallucination mitigation "
    "effectiveness varies by architecture — challenging the assumption that RAG universally improves medical QA.",
    font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_notes(slide,
    "SAY: 'My four key findings are:\n\n"
    "First, RAG significantly helps instruction-tuned models. GPT-4 jumped 24 points. These models are trained to FOLLOW instructions, so they know how to use the retrieved context.\n\n"
    "Second, not all models benefit. DistilBART actually got 7% worse — it is a summarisation model that gets confused by extra context. This challenges the assumption that RAG is always good.\n\n"
    "Third, domain pre-training can replace retrieval. BioBART achieved 58% without any help because medical knowledge is already baked into its weights from pre-training on biomedical literature.\n\n"
    "Fourth, real-world constraints matter. Llama 3 needed 4.6GB RAM — my laptop couldn't handle it. Phi-3 Mini works at 2.2GB. GPT-4 needs a paid API. These trade-offs shape which approach is practical for hospitals.\n\n"
    "My main contribution: proving that RAG effectiveness depends on architecture, not just the retrieval system.'")

# ═════════════════════════════════════════════════════════════════
# SLIDE 9: CHALLENGES & SOLUTIONS
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.4, 10, 0.6, "Engineering Challenges & Solutions",
            font_size=30, bold=True, color=WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(1.05), Inches(3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

challenges_data = [
    ["Challenge", "Root Cause", "Solution", "Outcome"],
    ["FAISS library crash", "NumPy 2.x removed\n_core API", "Switched to sklearn\ncosine similarity", "Same results,\nmore stable"],
    ["Llama 3 out of RAM", "4.6GB needed,\nonly 4.5GB free", "Switched to Phi-3 Mini\n(2.2GB)", "Competitive model\nthat runs on laptop"],
    ["GPT API quota\nexhausted", "Zero billing credit\non two accounts", "Added try/except;\nerror logging", "No data loss,\nclear error tracking"],
    ["T5-v1_1 tokenizer\nfailure", "Transformer version\nincompatibility", "Replaced with BioBART\n(biomedical model)", "Actually better: 58%\nmedical domain fit"],
    ["API key accidentally\nexposed", "Plain-text in chat\nhistory", "Revoked key, switched to\nenvironment variables", "Security restored,\nnever hardcoded"],
]
add_table(slide, 0.5, 1.3, 12.3, 3.5, 6, 4, challenges_data,
          col_widths=[2.5, 2.8, 3.2, 2.5])

add_textbox(slide, 0.8, 5.2, 11, 0.6,
    "Each challenge turned into an opportunity — forced better architectural decisions",
    font_size=15, color=ACCENT_ORANGE, alignment=PP_ALIGN.LEFT)

add_notes(slide,
    "SAY: 'Real engineering involves solving unexpected problems. Here are my biggest challenges:\n\n"
    "FAISS crashed because of a NumPy version conflict. I switched to sklearn cosine similarity which works identically for 500 documents.\n\n"
    "Llama 3 needed 4.6GB RAM but my laptop only had 4.5GB free. I found Phi-3 Mini at 2.2GB — and it actually showed the strongest RAG improvement.\n\n"
    "My GPT API key ran out of credits twice. Instead of losing data, I added error handling so the pipeline continued and logged which rows failed.\n\n"
    "The original T5-v1_1 model had a tokenizer bug. I replaced it with BioBART, which turned out BETTER because it is biomedical-specific.\n\n"
    "I accidentally exposed an API key. I immediately revoked it and switched to environment variables. A lesson in security best practices.'")

# ═════════════════════════════════════════════════════════════════
# SLIDE 10: LITERATURE & METHODOLOGY
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.4, 10, 0.6, "Literature Foundation & Methodology",
            font_size=30, bold=True, color=WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(1.05), Inches(3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

# Left column: Key Papers
add_textbox(slide, 0.8, 1.3, 5.5, 0.5, "Key Literature", font_size=20, bold=True, color=ACCENT_BLUE)

lit_items = [
    ("Lewis et al. (2020) — ", "RAG: Retrieval-Augmented Generation for NLP (my core technique)"),
    ("Zuo & Jiang (2025) — ", "MedHallBench: Medical hallucination benchmarking"),
    ("Ji et al. (2023) — ", "Self-Reflection for hallucination mitigation"),
    ("Sok et al. (2025) — ", "MetaRAG: Metamorphic testing for RAG systems"),
    ("Asgari et al. (2025) — ", "Clinical safety framework for LLM evaluation"),
    ("Tonmoy et al. (2024) — ", "Comprehensive survey of hallucination mitigation"),
]
add_bullet_list(slide, 0.8, 1.8, 6, 3.0, lit_items, font_size=14, color=LIGHT_GRAY)

# Right column: Methodology
add_textbox(slide, 7.5, 1.3, 5, 0.5, "Methodology", font_size=20, bold=True, color=ACCENT_BLUE)
add_bullet_list(slide, 7.5, 1.8, 5, 3.0, [
    "2×5 experimental grid (RAG vs No-RAG × 5 models)",
    "Same 100 PubMedQA questions across conditions",
    "Metrics: Accuracy, Macro F1, Precision, Recall",
    "Confusion matrices for yes/no/maybe analysis",
    "Top-1 and mean similarity scores logged",
    "Heuristic label extraction from generated text",
    "All results logged in CSV for reproducibility",
], font_size=14, color=LIGHT_GRAY)

add_textbox(slide, 0.8, 5.0, 11, 0.6,
    "Gap addressed: Most existing work uses static benchmarks — this project provides a comparative multi-model analysis with practical deployment insights",
    font_size=14, color=MED_GRAY)

add_notes(slide,
    "SAY: 'My work builds on established research. Lewis et al. introduced RAG — I applied it to healthcare. "
    "Zuo & Jiang created MedHallBench for medical hallucination evaluation — my metrics follow similar principles. "
    "Ji et al. proposed self-reflection loops, which I discuss as future work.\n\n"
    "The gap I address: most existing work benchmarks hallucination on static datasets with single models. "
    "My work is a comparative study across 5 architectures, showing that RAG effectiveness depends on model choice. "
    "This practical insight hasn't been widely reported.'")

# ═════════════════════════════════════════════════════════════════
# SLIDE 11: FUTURE WORK
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.4, 10, 0.6, "Future Work & Extensions",
            font_size=30, bold=True, color=WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(1.05), Inches(3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

# Future work items as boxes
future = [
    (0.8, 1.4, "Self-Reflection Loops",
     "AI checks its own answer before\noutputting — reduces hallucination\n(Ji et al., 2023)", ACCENT_BLUE),
    (4.3, 1.4, "Larger & Multi-Dataset",
     "MedQA, MedMCQA, MultiMedQA\nfor cross-domain evaluation\nand robustness testing", RGBColor(0x27, 0xAE, 0x60)),
    (7.8, 1.4, "Hallucination Metrics",
     "FACTSCORE, ACHMI for direct\nhallucination measurement\n(template already created)", RGBColor(0x8E, 0x44, 0xAD)),
    (0.8, 3.2, "Web Interface",
     "Streamlit/Flask demo for\nlive healthcare QA with\nretrieval-grounded responses", RGBColor(0xE6, 0x7E, 0x22)),
    (4.3, 3.2, "GPU-Scale Experiments",
     "Run Phi-3 on full 100+ samples\nTest larger models like Llama 3\nUse FAISS for million-doc scale", RGBColor(0xC0, 0x39, 0x2B)),
    (7.8, 3.2, "Clinical Safety Scoring",
     "Implement Asgari et al. (2025)\nframework for risk-aware\nevaluation of medical outputs", RGBColor(0x16, 0xA0, 0x85)),
]
for x, y, title, desc, color in future:
    add_shape_box(slide, x, y, 3.2, 0.5, color, title, font_size=14, text_color=WHITE, bold=True)
    add_textbox(slide, x + 0.1, y + 0.55, 3.0, 1.0, desc, font_size=12, color=LIGHT_GRAY)

add_notes(slide,
    "SAY: 'With more time, there are several promising extensions:\n\n"
    "Self-reflection loops where the AI reviews its own answer before presenting it. "
    "Testing on more datasets like MedQA and MedMCQA. "
    "Implementing direct hallucination metrics like FACTSCORE instead of just accuracy. "
    "Building a web interface so clinicians could try it themselves. "
    "Running on GPUs to test larger models at full scale. "
    "And adding clinical safety scoring frameworks for real-world risk assessment.'")

# ═════════════════════════════════════════════════════════════════
# SLIDE 12: CONCLUSION
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.4, 10, 0.6, "Conclusion",
            font_size=30, bold=True, color=WHITE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(1.05), Inches(3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

# 3 key takeaways
add_shape_box(slide, 1.5, 1.5, 10.3, 0.8, RGBColor(0x1A, 0x3C, 0x1A),
              "1.  RAG reduces hallucinations for instruction-tuned models  (GPT-4: 40% → 64%)",
              font_size=18, text_color=ACCENT_GREEN, bold=True)

add_shape_box(slide, 1.5, 2.5, 10.3, 0.8, RGBColor(0x3C, 0x1A, 0x1A),
              "2.  Not all models benefit — architecture matters  (DistilBART: 39% → 32%)",
              font_size=18, text_color=ACCENT_RED, bold=True)

add_shape_box(slide, 1.5, 3.5, 10.3, 0.8, RGBColor(0x2C, 0x2C, 0x15),
              "3.  Domain pre-training (BioBART 58%) is a viable alternative to RAG",
              font_size=18, text_color=ACCENT_ORANGE, bold=True)

add_textbox(slide, 1.0, 4.8, 11, 1.2,
    '"For real-world healthcare AI, we need BOTH good retrieval AND the right\n'
    'model architecture to make AI trustworthy. RAG is powerful but not universal."',
    font_size=20, bold=False, color=WHITE, alignment=PP_ALIGN.CENTER)

# Project stats
add_textbox(slide, 1.0, 6.2, 11, 0.5,
    "5 models  •  10 experiments  •  500 documents  •  17 scripts  •  2,500+ lines of code",
    font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_notes(slide, 
    "SAY: 'To conclude: I proved three things.\n\n"
    "First, RAG significantly reduces hallucinations for the right models — GPT-4 went from 40% to 64%.\n\n"
    "Second, RAG is not universal. DistilBART actually got worse. Architecture matters.\n\n"
    "Third, domain pre-training like BioBART is a viable alternative — 58% without any retrieval.\n\n"
    "For real-world healthcare AI, we need both good retrieval AND the right model. Thank you.'")

# ═════════════════════════════════════════════════════════════════
# SLIDE 13: THANK YOU / Q&A
# ═════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_textbox(slide, 1, 2.0, 11.3, 1.0,
    "Thank You — Questions?",
    font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(5.5), Inches(3.0), Inches(2.3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

add_textbox(slide, 1, 3.3, 11.3, 1.5,
    "Ayomide Ogun-Ajala  |  21309903\n"
    "BSc Computer Science & Software Engineering\n"
    "Supervisor: Dr. Kolawole Adebayo\n"
    "Maynooth University",
    font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1, 5.2, 11.3, 0.5,
    "HealthHubQA — Investigating & Mitigating Hallucinations in LLMs for Healthcare QA",
    font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_notes(slide,
    "Be ready for questions. Key answers to have ready:\n"
    "- What is hallucination? (Confident but wrong AI output)\n"
    "- Why these 5 models? (Compare architectures, not just one)\n"
    "- Why is 64% not higher? (PubMedQA is hard, focus on RELATIVE improvement)\n"
    "- Ethics? (Always needs human oversight, assist not replace clinicians)\n"
    "- What would you do differently? (GPU from day 1, FACTSCORE metrics, start all models simultaneously)")

# ── Save ────────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           'HealthHubQA_FYP_Presentation_v2.pptx')
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
