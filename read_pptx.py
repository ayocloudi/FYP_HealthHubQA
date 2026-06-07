from pptx import Presentation

prs = Presentation('HealthHubQA_FYP_Presentation.pptx')

for i, slide in enumerate(prs.slides):
    layout_name = slide.slide_layout.name if slide.slide_layout else 'Unknown'
    print(f'--- Slide {i+1} (Layout: {layout_name}) ---')
    for shape in slide.shapes:
        print(f'  Shape: {shape.shape_type}, Name: {shape.name}, Pos: ({shape.left}, {shape.top}), Size: ({shape.width}, {shape.height})')
        if shape.has_text_frame:
            for p_idx, para in enumerate(shape.text_frame.paragraphs):
                text = para.text.strip()
                if text:
                    print(f'    Para {p_idx}: "{text[:200]}"')
        if shape.has_table:
            table = shape.table
            print(f'    Table: {len(table.rows)} rows x {len(table.columns)} cols')
            for r_idx, row in enumerate(table.rows):
                cells = [cell.text.strip() for cell in row.cells]
                print(f'      Row {r_idx}: {cells}')
    print()

print(f'Total slides: {len(prs.slides)}')
